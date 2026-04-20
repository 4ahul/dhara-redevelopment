"""
Unified Document Extractor
Handles PDF (text + scanned), DOCX, DOC, XLSX, TIF/TIFF, JPEG, PPTX
with multilingual OCR (English + Marathi + Hindi).
"""

import os
import sys
import time
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from concurrent.futures import ThreadPoolExecutor, as_completed

# Add parent to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))


@dataclass
class ExtractedPage:
    """A single page of extracted content with metadata."""

    text: str
    page_num: int
    method: str  # "pypdf", "ocr", "docx", "xlsx", "pptx", "image"
    language: str = "unknown"
    confidence: float = 0.0


@dataclass
class ExtractedDocument:
    """Complete extracted document with metadata."""

    filepath: str
    filename: str
    relative_path: str
    file_hash: str
    file_size: int
    pages: List[ExtractedPage] = field(default_factory=list)
    doc_type: str = "other"
    total_chars: int = 0
    extraction_method: str = ""
    errors: List[str] = field(default_factory=list)

    @property
    def full_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages if p.text)


# ---------------------------------------------------------------------------
# OCR Engines
# ---------------------------------------------------------------------------

_EASYOCR_READER = None
_TESSERACT_AVAILABLE = False


def _get_easyocr_reader():
    """Lazy-init EasyOCR reader with English + Marathi + Hindi."""
    global _EASYOCR_READER
    if _EASYOCR_READER is None:
        try:
            import easyocr

            print("  [OCR] Initializing EasyOCR (en, mr, hi)...")
            _EASYOCR_READER = easyocr.Reader(
                ["en", "mr", "hi"], gpu=False, verbose=False
            )
            print("  [OCR] EasyOCR ready")
        except Exception as e:
            print(f"  [OCR] EasyOCR not available: {e}")
    return _EASYOCR_READER


def _check_tesseract():
    """Check if Tesseract is available."""
    global _TESSERACT_AVAILABLE
    try:
        import pytesseract

        # Try to run tesseract
        pytesseract.get_tesseract_version()
        _TESSERACT_AVAILABLE = True
    except Exception:
        _TESSERACT_AVAILABLE = False
    return _TESSERACT_AVAILABLE


def ocr_image_easyocr(image_path: str) -> Optional[str]:
    """Run EasyOCR on a single image."""
    reader = _get_easyocr_reader()
    if reader is None:
        return None
    try:
        results = reader.readtext(image_path, detail=0, paragraph=True)
        return "\n".join(results).strip() if results else None
    except Exception:
        return None


def ocr_image_tesseract(image_path: str) -> Optional[str]:
    """Run Tesseract on a single image with Devanagari support."""
    if not _check_tesseract():
        return None
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(image_path)
        # eng+mar+hin for English+Marathi+Hindi
        text = pytesseract.image_to_string(img, lang="eng+mar+hin")
        return text.strip() if text else None
    except Exception:
        return None


def ocr_pil_image_easyocr(img) -> Optional[str]:
    """Run EasyOCR on a PIL Image object."""
    reader = _get_easyocr_reader()
    if reader is None:
        return None
    try:
        import numpy as np

        img_array = np.array(img)
        results = reader.readtext(img_array, detail=0, paragraph=True)
        return "\n".join(results).strip() if results else None
    except Exception:
        return None


# ---------------------------------------------------------------------------
# PDF Extraction
# ---------------------------------------------------------------------------


def extract_pdf_text(filepath: Path) -> List[ExtractedPage]:
    """Extract text from PDF using pypdf (text-based PDFs)."""
    pages = []
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(filepath))
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = text.strip()
            if text:
                pages.append(
                    ExtractedPage(
                        text=text,
                        page_num=i + 1,
                        method="pypdf",
                    )
                )
    except Exception as e:
        pass
    return pages


def extract_pdf_ocr(filepath: Path) -> List[ExtractedPage]:
    """Extract text from scanned PDF using OCR (pdf2image + EasyOCR/Tesseract)."""
    pages = []

    # Try pdf2image first
    try:
        from pdf2image import convert_from_path

        images = convert_from_path(str(filepath), dpi=200, first_page=1, last_page=100)
    except Exception:
        # Fallback: try converting individual pages via PIL for TIFF-based PDFs
        return []

    for i, img in enumerate(images):
        text = None

        # Try EasyOCR first (better for mixed Devanagari/Latin)
        text = ocr_pil_image_easyocr(img)
        method = "ocr_easyocr"

        # Fallback to Tesseract
        if not text or len(text) < 20:
            try:
                import pytesseract

                text_tess = pytesseract.image_to_string(img, lang="eng+mar+hin")
                if text_tess and len(text_tess.strip()) > len(text or ""):
                    text = text_tess.strip()
                    method = "ocr_tesseract"
            except Exception:
                pass

        if text and len(text) > 10:
            pages.append(
                ExtractedPage(
                    text=text,
                    page_num=i + 1,
                    method=method,
                )
            )

    return pages


def extract_pdf(filepath: Path) -> Tuple[List[ExtractedPage], str]:
    """Extract PDF - try text first, fall back to OCR."""
    pages = extract_pdf_text(filepath)
    if pages and sum(len(p.text) for p in pages) > 200:
        return pages, "pypdf"

    # Text extraction insufficient, try OCR
    ocr_pages = extract_pdf_ocr(filepath)
    if ocr_pages:
        return ocr_pages, "ocr"

    # Return whatever we got
    return pages, "pypdf" if pages else "none"


# ---------------------------------------------------------------------------
# Image Extraction (JPEG, TIFF, PNG)
# ---------------------------------------------------------------------------


def extract_image(filepath: Path) -> Tuple[List[ExtractedPage], str]:
    """Extract text from image files using OCR."""
    text = ocr_image_easyocr(str(filepath))
    method = "ocr_easyocr"

    if not text or len(text) < 20:
        text_tess = ocr_image_tesseract(str(filepath))
        if text_tess and len(text_tess) > len(text or ""):
            text = text_tess
            method = "ocr_tesseract"

    if text and len(text) > 10:
        return [ExtractedPage(text=text, page_num=1, method=method)], method
    return [], "none"


# ---------------------------------------------------------------------------
# DOCX Extraction
# ---------------------------------------------------------------------------


def extract_docx(filepath: Path) -> Tuple[List[ExtractedPage], str]:
    """Extract text from DOCX files."""
    try:
        import docx

        doc = docx.Document(str(filepath))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)

        # Also extract text from tables
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    table_text.append(row_text)

        if table_text:
            text += "\n\n" + "\n".join(table_text)

        if text.strip():
            return [ExtractedPage(text=text.strip(), page_num=1, method="docx")], "docx"
    except Exception as e:
        pass
    return [], "none"


# ---------------------------------------------------------------------------
# DOC Extraction (legacy format)
# ---------------------------------------------------------------------------


def extract_doc(filepath: Path) -> Tuple[List[ExtractedPage], str]:
    """Extract text from legacy .doc files."""
    # Try python-docx2txt or antiword
    try:
        import subprocess

        result = subprocess.run(
            ["catdoc", str(filepath)], capture_output=True, text=True, timeout=30
        )
        text = result.stdout.strip()
        if text and len(text) > 20:
            return [ExtractedPage(text=text, page_num=1, method="catdoc")], "catdoc"
    except Exception:
        pass

    # Fallback: try reading as raw text
    try:
        with open(filepath, "rb") as f:
            raw = f.read()
        # Extract printable ASCII and Devanagari sequences
        text = raw.decode("utf-8", errors="ignore").strip()
        if len(text) > 50:
            return [ExtractedPage(text=text, page_num=1, method="raw_doc")], "raw"
    except Exception:
        pass

    return [], "none"


# ---------------------------------------------------------------------------
# XLSX Extraction
# ---------------------------------------------------------------------------


def extract_xlsx(filepath: Path) -> Tuple[List[ExtractedPage], str]:
    """Extract text from XLSX files."""
    try:
        import openpyxl

        wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
        pages = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    rows.append(row_text)

            if rows:
                text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
                pages.append(ExtractedPage(text=text, page_num=1, method="xlsx"))

        wb.close()
        if pages:
            return pages, "xlsx"
    except Exception:
        pass
    return [], "none"


# ---------------------------------------------------------------------------
# PPTX Extraction
# ---------------------------------------------------------------------------


def extract_pptx(filepath: Path) -> Tuple[List[ExtractedPage], str]:
    """Extract text from PPTX files."""
    try:
        from pptx import Presentation

        prs = Presentation(str(filepath))
        text_parts = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_text))

        if text_parts:
            text = "\n\n".join(text_parts)
            return [ExtractedPage(text=text, page_num=1, method="pptx")], "pptx"
    except Exception:
        pass
    return [], "none"


# ---------------------------------------------------------------------------
# Main extraction dispatcher
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".doc",
    ".xlsx",
    ".xls",
    ".pptx",
    ".tif",
    ".tiff",
    ".jpg",
    ".jpeg",
    ".jpe",
    ".png",
    ".bmp",
    ".txt",
    ".csv",
}

SKIP_FILES = {
    "Thumbs.db",
    "desktop.ini",
    ".DS_Store",
}

SKIP_EXTENSIONS = {".tmp", ".bak", ".db", ".lnk", ".ini"}


def compute_file_hash(filepath: Path) -> str:
    """Compute MD5 hash of file for deduplication."""
    h = hashlib.md5()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def extract_document(filepath: Path, docs_root: Path) -> ExtractedDocument:
    """Extract a single document. Returns ExtractedDocument."""
    filename = filepath.name
    ext = filepath.suffix.lower()

    doc = ExtractedDocument(
        filepath=str(filepath),
        filename=filename,
        relative_path=str(filepath.relative_to(docs_root)),
        file_hash=compute_file_hash(filepath),
        file_size=filepath.stat().st_size,
    )

    if filename in SKIP_FILES or ext in SKIP_EXTENSIONS:
        doc.errors.append("Skipped file")
        return doc

    if ext not in SUPPORTED_EXTENSIONS:
        doc.errors.append(f"Unsupported extension: {ext}")
        return doc

    try:
        if ext == ".pdf":
            pages, method = extract_pdf(filepath)
        elif ext == ".docx":
            pages, method = extract_docx(filepath)
        elif ext == ".doc":
            pages, method = extract_doc(filepath)
        elif ext in {".xlsx", ".xls"}:
            pages, method = extract_xlsx(filepath)
        elif ext == ".pptx":
            pages, method = extract_pptx(filepath)
        elif ext in {".tif", ".tiff", ".jpg", ".jpeg", ".jpe", ".png", ".bmp"}:
            pages, method = extract_image(filepath)
        elif ext in {".txt", ".csv"}:
            try:
                text = filepath.read_text(encoding="utf-8", errors="ignore")
                pages = [ExtractedPage(text=text.strip(), page_num=1, method="txt")]
                method = "txt"
            except Exception:
                pages, method = [], "none"
        else:
            pages, method = [], "none"

        doc.pages = pages
        doc.extraction_method = method
        doc.total_chars = sum(len(p.text) for p in pages)

    except Exception as e:
        doc.errors.append(str(e))

    return doc


def find_all_documents(docs_dir: Path) -> List[Path]:
    """Find all supported documents recursively."""
    files = []
    for root, dirs, filenames in os.walk(docs_dir):
        # Skip hidden directories and cache
        dirs[:] = [d for d in dirs if not d.startswith((".", "__"))]
        for f in filenames:
            if f.startswith((".", "~$")):
                continue
            ext = os.path.splitext(f)[1].lower()
            if ext in SUPPORTED_EXTENSIONS:
                files.append(Path(root) / f)
    return sorted(files)


def extract_all_documents(
    docs_dir: str = "data/docs",
    max_workers: int = 4,
    verbose: bool = True,
) -> List[ExtractedDocument]:
    """
    Extract all documents from a directory.
    Returns list of ExtractedDocument with text content.
    """
    docs_root = Path(docs_dir)
    files = find_all_documents(docs_root)

    if verbose:
        total_size = sum(f.stat().st_size for f in files) / 1024 / 1024
        print(f"[EXTRACT] Found {len(files)} documents ({total_size:.1f} MB)")

    results = []
    t0 = time.time()

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(extract_document, f, docs_root): f for f in files}
        for i, future in enumerate(as_completed(futures)):
            try:
                doc = future.result()
                results.append(doc)
                if verbose and (i + 1) % 50 == 0:
                    elapsed = time.time() - t0
                    rate = (i + 1) / elapsed
                    total_chars = sum(d.total_chars for d in results)
                    print(
                        f"  Progress: {i + 1}/{len(files)} ({rate:.1f} files/s, {total_chars / 1000:.0f}K chars)"
                    )
            except Exception as e:
                if verbose:
                    f = futures[future]
                    print(f"  Error processing {f.name}: {e}")

    if verbose:
        elapsed = time.time() - t0
        success = len([d for d in results if d.total_chars > 0])
        total_chars = sum(d.total_chars for d in results)
        print(f"\n[EXTRACT] Done in {elapsed:.1f}s")
        print(f"  Successfully extracted: {success}/{len(results)} documents")
        print(f"  Total characters: {total_chars:,}")
        methods = {}
        for d in results:
            methods[d.extraction_method] = methods.get(d.extraction_method, 0) + 1
        for m, c in sorted(methods.items(), key=lambda x: -x[1]):
            print(f"  {m}: {c} docs")

    return results


if __name__ == "__main__":
    docs_dir = sys.argv[1] if len(sys.argv) > 1 else "data/docs"
    results = extract_all_documents(docs_dir)
